from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_portfolio():
    prs = Presentation()
    
    # 세로형(A4) 설정 (8.27 x 11.69 inches)
    prs.slide_width = Inches(8.27)
    prs.slide_height = Inches(11.69)

    projects = [
        {
            "title": "MyWebPage",
            "subtitle": "Django & Docker 기반 풀스택 플랫폼",
            "purpose": "Django 아키텍처 이해 및 Docker 배포 환경 구축",
            "tech": "Django, Python, Docker, Nginx, Certbot",
            "features": ["Full-stack SSR 구현", "Containerization 기반 배포", "HTTPS 보안 적용"],
            "review": "웹 서비스 전체 사이클 이해 및 인프라 최적화 역량 습득",
            "period": "2023.01 - 2023.06"
        },
        {
            "title": "VesselRouteVisualizer",
            "subtitle": "해상 노선 시각화 및 데이터 분석 플랫폼",
            "purpose": "해상 물류 데이터의 시각적 가독성 및 관리 효율 증대",
            "tech": "FastAPI, SQLAlchemy, Leaflet.js, Python",
            "features": ["Interactive Map 시각화", "ETL 데이터 적재 파이프라인", "비동기 API 서버 구축"],
            "review": "GIS 데이터 최적화 및 고성능 API 설계 능력 향상",
            "period": "2023.08 - 2024.02"
        },
        {
            "title": "pyacet",
            "subtitle": "특허 정보 수집 및 데이터 정제 자동화 툴",
            "purpose": "연구용 특허 데이터 수집 자동화 및 데이터 정제",
            "tech": "Python, Selenium, Pandas, BeautifulSoup",
            "features": ["동적 웹 크롤링 자동화", "데이터 클렌징 및 정형화", "Excel/CSV 리포트 생성"],
            "review": "데이터 엔지니어링 파이프라인 구축 역량 강화",
            "period": "2023.11 - 2024.01"
        }
    ]

    for p in projects:
        slide = prs.slides.add_slide(prs.slide_layouts[6]) # 빈 슬라이드
        
        # 제목 배경 (간단한 상자)
        left = Inches(0.5); top = Inches(0.5); width = Inches(7.27); height = Inches(1)
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_box.text = p["title"]
        title_box.text_frame.paragraphs[0].font.size = Pt(32)
        title_box.text_frame.paragraphs[0].font.bold = True
        
        # 내용 구성 (목적, 스택, 기능 등)
        body_box = slide.shapes.add_textbox(left, Inches(1.8), width, Inches(8))
        tf = body_box.text_frame
        tf.word_wrap = True
        
        def add_para(text, size, bold=False):
            para = tf.add_paragraph()
            para.text = text
            para.font.size = Pt(size)
            para.font.bold = bold
            para.space_after = Pt(10)

        add_para(f"■ {p['subtitle']}", 18, True)
        add_para(f"▶ 목적: {p['purpose']}", 14)
        add_para(f"▶ 기술 스택: {p['tech']}", 14, True)
        add_para("▶ 핵심 기능:", 14)
        for f in p["features"]:
            add_para(f"  - {f}", 12)
        add_para(f"▶ 배운점 및 후기: {p['review']}", 14)
        add_para(f"▶ 프로젝트 기간: {p['period']}", 12)

    prs.save('Portfolio_LinfoKR_Vertical.pptx')
    print("PPTX 파일이 성공적으로 생성되었습니다!")

if __name__ == "__main__":
    create_portfolio()